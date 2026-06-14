import logging
from pptx import Presentation

logger = logging.getLogger(__name__)

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PowerPoint file"""
    try:
        prs = Presentation(file_path)
        text = ""
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"Slide {slide_num}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            text += "\n"
        
        return text.strip()
    except Exception as e:
        logger.error(f"Error extracting text from PPTX: {e}")
        raise
